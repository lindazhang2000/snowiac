"""Generate the SnowIaC pre-demo PowerPoint deck (docs/SnowIaC_Demo.pptx)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs" / "SnowIaC_Demo.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

# Brand palette
NAVY = RGBColor(0x0B, 0x2A, 0x4A)
BLUE = RGBColor(0x1F, 0x6F, 0xEB)
TEAL = RGBColor(0x15, 0xA0, 0x8C)
GOLD = RGBColor(0xE6, 0xA4, 0x23)
GRAY = RGBColor(0x55, 0x5B, 0x66)
LIGHT = RGBColor(0xF3, 0xF6, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1F, 0x2C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height

BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    return s


def notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.clear()
    paras = text.strip().split("\n\n")
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = para.strip()
        p.font.size = Pt(12)


def header_bar(s, title, subtitle=None):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.15), SW - Inches(1), Inches(0.7))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = RGBColor(0xCF, 0xDA, 0xE8)


def textbox(s, x, y, w, h, lines, font_size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    if isinstance(lines, str):
        lines = [lines]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if isinstance(line, tuple):
            text, opts = line
        else:
            text, opts = line, {}
        p.text = text
        p.font.size = Pt(opts.get("size", font_size))
        p.font.bold = opts.get("bold", bold)
        p.font.color.rgb = opts.get("color", color)
        if opts.get("space_after"):
            p.space_after = Pt(opts["space_after"])
    return tb


def box(s, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = s.shapes.add_shape(shape, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def label(shp, text, color=WHITE, size=14, bold=True, subtitle=None, sub_size=10, sub_color=None):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.margin_left = tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if subtitle:
        sp = tf.add_paragraph()
        sp.alignment = PP_ALIGN.CENTER
        sp.text = subtitle
        sp.font.size = Pt(sub_size)
        sp.font.bold = False
        sp.font.color.rgb = sub_color or color


def arrow(s, x1, y1, x2, y2, color=NAVY, weight=2.25):
    line = s.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # Add arrowhead via XML
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        from lxml import etree
        tail = etree.SubElement(ln, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("h", "med")
    return line


def footer(s, text):
    tb = s.shapes.add_textbox(Inches(0.4), SH - Inches(0.4), SW - Inches(0.8), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = GRAY


# ─── SLIDE 1: Title ───────────────────────────────────────────────────────────
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.6), SW, Inches(0.12))
accent.line.fill.background()
accent.fill.solid()
accent.fill.fore_color.rgb = GOLD

textbox(s, Inches(0.8), Inches(1.7), SW - Inches(1.6), Inches(1.2),
        "SnowIaC", font_size=72, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(2.9), SW - Inches(1.6), Inches(0.8),
        "ServiceNow → AI Agents → Infrastructure as Code → Azure",
        font_size=24, color=RGBColor(0xCF, 0xDA, 0xE8))
textbox(s, Inches(0.8), Inches(3.7), SW - Inches(1.6), Inches(0.6),
        "Auto-pilot the ITSM-to-cloud workflow with Microsoft Agent Framework + Azure Foundry",
        font_size=16, color=RGBColor(0xA9, 0xBC, 0xD2))
textbox(s, Inches(0.8), Inches(6.0), SW - Inches(1.6), Inches(0.5),
        "Demo • Live against real ServiceNow, real GitHub, real Azure",
        font_size=14, color=GOLD, bold=True)

notes(s, """
Welcome. In the next ~15 minutes I'll show you SnowIaC — an end-to-end pipeline that turns a ServiceNow change request into a deployed, verified Azure resource with zero manual steps after the human PR approval.

The name is the workflow: SNOW (ServiceNow) → IaC (Infrastructure as Code). Built on Microsoft Agent Framework and Azure Foundry.

Everything you'll see is live. Real ServiceNow instance, real GitHub repository, real Azure resources. No mocks, no replays. Total runtime end-to-end is about 3-4 minutes.
""")

# ─── SLIDE 2: The Problem ─────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "The Problem", "Infrastructure change requests are slow, manual, and error-prone")

problems = [
    ("Tickets pile up", "Cloud admins triage every RITM by hand — even routine disk resizes."),
    ("Context-switching tax", "SNOW → wiki → Terraform repo → PR review → Azure portal → SNOW close."),
    ("Inconsistent code", "Hand-written Terraform drifts from standards; reviewers re-litigate every PR."),
    ("Slow verification", "Did the change actually land? Someone has to log in and check."),
    ("No audit trail", "Approvals, code, deploy, verify, closure spread across 5 systems."),
]
x0, y0 = Inches(0.6), Inches(1.4)
card_w, card_h = Inches(6.0), Inches(1.0)
for i, (title, body) in enumerate(problems):
    row, col = divmod(i, 2)
    x = x0 + col * (card_w + Inches(0.2))
    y = y0 + row * (card_h + Inches(0.25))
    b = box(s, x, y, card_w, card_h, LIGHT, line=RGBColor(0xD0, 0xD7, 0xE0))
    label(b, title, color=NAVY, size=16, subtitle=body, sub_size=12, sub_color=GRAY)

textbox(s, Inches(0.6), Inches(6.3), SW - Inches(1.2), Inches(0.6),
        "→ Result: hours-to-days per ticket, with humans doing work agents can do reliably.",
        font_size=16, color=BLUE, bold=True)

notes(s, """
Let's frame the problem. In any organization running cloud at scale, infrastructure change requests come in through ITSM — ServiceNow being the dominant system.

Five things go wrong every time:

1. Tickets pile up. Even a simple disk IOPS increase has to be triaged by a cloud admin.

2. Context-switching tax. The admin flips between ServiceNow to read the ticket, an internal wiki for the standard, a Terraform repo to write the code, GitHub to open a PR, Azure portal to verify, and back to ServiceNow to close. Six systems, one ticket.

3. Hand-written Terraform drifts. Every engineer writes it slightly differently and PR reviewers re-litigate the same patterns weekly.

4. Verification is manual. Did the change actually land with the right values? Someone logs in and eyeballs it.

5. The audit trail is split across five systems. Auditors hate this.

Net result: routine changes take hours to days, and humans are doing work that agents can do reliably and faster.
""")

# ─── SLIDE 3: What SnowIaC Does ───────────────────────────────────────────────
s = add_slide()
header_bar(s, "What SnowIaC Does", "One pipeline. Five agents. Zero ticket-shuffling.")

bullets = [
    ("Ingests", "Polls ServiceNow for new RITMs (or accepts /tickets/intake POSTs)."),
    ("Classifies", "Intake agent uses Foundry gpt-5.4 to pick the request kind and validate fields."),
    ("Generates", "Code-gen agent emits Terraform from Jinja2 templates with strict undefined checks."),
    ("Proposes", "Opens a GitHub PR with the Terraform diff and human-readable summary."),
    ("Deploys", "On merge, GH Actions runs terraform apply against Azure."),
    ("Verifies", "Verification agent reads the live resource and proves spec compliance."),
    ("Closes", "Closure agent patches the SNOW ticket to Closed Complete with evidence."),
]
y = Inches(1.3)
for tag, body in bullets:
    chip = box(s, Inches(0.7), y, Inches(1.5), Inches(0.55), BLUE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label(chip, tag, color=WHITE, size=14)
    textbox(s, Inches(2.4), y, Inches(10.3), Inches(0.55),
            body, font_size=15, color=DARK)
    y += Inches(0.65)

textbox(s, Inches(0.7), Inches(6.5), SW - Inches(1.4), Inches(0.5),
        "Human in the loop ONLY at PR review — every other step is autonomous and auditable.",
        font_size=14, color=TEAL, bold=True)

notes(s, """
SnowIaC takes that whole flow and collapses it into one pipeline with five specialized agents.

Ingest: polls ServiceNow or accepts a webhook. Routine plumbing.

Classify: the Intake agent uses Foundry gpt-5.4 to figure out what kind of request this is and whether the required fields are present. If something's missing, it tells ServiceNow and stops.

Generate: the CodeGen agent takes the structured ticket data and renders Terraform from Jinja2 templates with StrictUndefined — meaning if any required variable is missing, generation fails loudly rather than producing bad code.

Propose: opens a GitHub PR with the diff and a human-readable summary so a reviewer knows exactly what's changing.

Deploy: when the PR is merged, GitHub Actions runs terraform apply with OIDC — no long-lived credentials.

Verify: the Verification agent reads the live Azure resource and compares it to the spec from the original ticket.

Close: finally the Closure agent updates ServiceNow to Closed Complete and attaches the evidence.

The key design point: the human stays in the loop, but only at the PR review step. Every other transition is autonomous, observable, and fully audited.
""")

# ─── SLIDE 4: High-Level Architecture (diagram) ───────────────────────────────
s = add_slide()
header_bar(s, "High-Level Architecture", "ServiceNow → Foundry Agents → GitHub → Azure")

# Three swim lanes background
lane_y = Inches(1.2)
lane_h = Inches(5.4)
lane_w = (SW - Inches(0.8)) / 3
lane_titles = [("ServiceNow", GOLD), ("SnowIaC App (Container App + Foundry Agents)", BLUE), ("Azure / GitHub", TEAL)]
for i, (t, c) in enumerate(lane_titles):
    lx = Inches(0.4) + i * lane_w
    band = box(s, lx, lane_y, lane_w - Inches(0.1), Inches(0.4), c, shape=MSO_SHAPE.RECTANGLE)
    label(band, t, color=WHITE, size=13)
    body = box(s, lx, lane_y + Inches(0.4), lane_w - Inches(0.1), lane_h - Inches(0.4),
               LIGHT, line=RGBColor(0xD0, 0xD7, 0xE0), shape=MSO_SHAPE.RECTANGLE)

# Lane 1: ServiceNow
snow = box(s, Inches(0.6), Inches(2.0), Inches(3.7), Inches(1.0), GOLD)
label(snow, "ServiceNow RITM", color=WHITE, size=15,
      subtitle="Catalog item: Azure Cloud Infrastructure Change",
      sub_size=10, sub_color=WHITE)

snow_close = box(s, Inches(0.6), Inches(5.4), Inches(3.7), Inches(1.0), GRAY)
label(snow_close, "RITM → Closed Complete", color=WHITE, size=14,
      subtitle="With deploy + verification evidence",
      sub_size=10, sub_color=WHITE)

# Lane 2: Agents
agents = [
    ("Intake", "Classify + validate", Inches(2.0)),
    ("CodeGen", "Render Terraform", Inches(2.9)),
    ("PR", "Open GitHub PR", Inches(3.8)),
    ("Verify", "Read live resource", Inches(4.7)),
    ("Close", "Patch SNOW", Inches(5.6)),
]
ax = Inches(4.8)
for name, sub, y in agents:
    b = box(s, ax, y, Inches(3.6), Inches(0.7), BLUE)
    label(b, name + "  Agent", color=WHITE, size=13, subtitle=sub, sub_size=10, sub_color=RGBColor(0xCF, 0xDA, 0xE8))

# Lane 3: Azure / GitHub
gh = box(s, Inches(9.2), Inches(2.0), Inches(3.7), Inches(1.0), DARK)
label(gh, "GitHub PR + Actions", color=WHITE, size=15,
      subtitle="terraform-apply workflow on merge",
      sub_size=10, sub_color=RGBColor(0xCF, 0xDA, 0xE8))

azure = box(s, Inches(9.2), Inches(3.4), Inches(3.7), Inches(1.0), TEAL)
label(azure, "Azure Resources", color=WHITE, size=15,
      subtitle="Managed disks, VMs, etc. (target RG)",
      sub_size=10, sub_color=WHITE)

foundry = box(s, Inches(9.2), Inches(4.8), Inches(3.7), Inches(1.0), NAVY)
label(foundry, "Azure Foundry", color=WHITE, size=15,
      subtitle="gpt-5.4 deployment • UAMI auth",
      sub_size=10, sub_color=RGBColor(0xCF, 0xDA, 0xE8))

# Arrows: SNOW → Intake
arrow(s, Inches(4.3), Inches(2.5), Inches(4.8), Inches(2.35))
# CodeGen → GH
arrow(s, Inches(8.4), Inches(3.25), Inches(9.2), Inches(2.5))
# Agents ↔ Foundry
arrow(s, Inches(8.4), Inches(4.4), Inches(9.2), Inches(5.3), color=NAVY)
arrow(s, Inches(9.2), Inches(5.0), Inches(8.4), Inches(4.6), color=NAVY)
# GH Actions → Azure
arrow(s, Inches(11.05), Inches(3.0), Inches(11.05), Inches(3.4), color=TEAL)
# Verify reads Azure
arrow(s, Inches(9.2), Inches(3.9), Inches(8.4), Inches(5.05), color=TEAL)
# Close → SNOW closed
arrow(s, Inches(4.8), Inches(5.95), Inches(4.3), Inches(5.9), color=GRAY)

footer(s, "Webhook back from GH Actions (HMAC-signed) triggers the Verify → Close cascade.")

notes(s, """
This is the high-level architecture. Three swim lanes: ServiceNow on the left, the SnowIaC app in the middle running on Azure Container Apps, and the downstream systems — GitHub and Azure — on the right.

The flow reads left-to-right, top-to-bottom:

A RITM is created in ServiceNow. Either polled or pushed via the catalog item we built.

The Intake agent picks it up, classifies it with Foundry, and validates required fields.

The CodeGen agent renders Terraform and the PR agent opens a pull request in GitHub.

A human merges. That triggers the terraform-apply GitHub Actions workflow, which applies against the target Azure resource group.

GitHub Actions sends an HMAC-signed webhook back to SnowIaC saying "the apply succeeded." That fires the Verify agent, which reads the live resource and confirms it matches the spec.

Finally the Close agent patches the RITM to Closed Complete with the apply log and verification evidence.

All five agents authenticate to Foundry using a user-assigned managed identity. No API keys anywhere. Everything is auditable.
""")

# ─── SLIDE 5: Tech Stack ──────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Tech Stack", "Cloud-native, fully managed, fully reproducible")

stacks = [
    ("Agents", BLUE, [
        "Microsoft Agent Framework 1.0.0rc6",
        "Azure Foundry • gpt-5.4 deployment",
        "ManagedIdentityCredential (UAMI)",
    ]),
    ("App", TEAL, [
        "FastAPI + Pydantic v2",
        "Postgres ticket store (Azure Flexible Server)",
        "Jinja2 (StrictUndefined) Terraform templates",
    ]),
    ("Infra", NAVY, [
        "Azure Container Apps + ACR",
        "Key Vault for secrets (RBAC)",
        "Terraform (azurerm ~> 4.0) — all in repo",
    ]),
    ("Pipeline", GOLD, [
        "GitHub Actions (OIDC, no long-lived creds)",
        "HMAC-signed deploy webhooks",
        "ServiceNow REST API",
    ]),
]
col_w = Inches(3.0)
gap = Inches(0.15)
x0 = Inches(0.5)
y = Inches(1.4)
for i, (title, color, items) in enumerate(stacks):
    x = x0 + i * (col_w + gap)
    head = box(s, x, y, col_w, Inches(0.55), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    label(head, title, color=WHITE, size=18)
    bod = box(s, x, y + Inches(0.65), col_w, Inches(4.5), LIGHT, line=RGBColor(0xD0, 0xD7, 0xE0))
    tf = bod.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)

textbox(s, Inches(0.5), Inches(6.5), SW - Inches(1), Inches(0.5),
        "Everything is in Terraform — including the Foundry role assignments. Reproducible from a clean subscription.",
        font_size=13, color=GRAY, bold=False)

notes(s, """
Quick tour of what's under the hood.

Agents: Microsoft Agent Framework version 1.0.0rc6 driving five Foundry agents on a gpt-5.4 deployment. Authentication is ManagedIdentityCredential — the Container App's user-assigned managed identity has the Azure AI User role on the Foundry account.

App: FastAPI with Pydantic v2 for type-safe ticket models. Postgres Flexible Server for persistence — durable ticket state survives container restarts. Jinja2 templates with StrictUndefined so we fail fast on missing data instead of generating broken Terraform.

Infra: Azure Container Apps for the SnowIaC server itself, Azure Container Registry for the image, Key Vault with RBAC for all secrets — HMAC, ServiceNow password, GitHub token, database connection string. Everything declared in Terraform using azurerm provider version 4.

Pipeline: GitHub Actions using OIDC federation — no long-lived service-principal secrets. HMAC-signed webhooks back from GitHub so SnowIaC knows the deploy ran. ServiceNow REST API on the upstream side.

Key point at the bottom: the entire stack — including the Foundry role assignments and the Postgres database — is in Terraform. From a clean subscription, you could redeploy this whole thing in 10 minutes.
""")

# ─── SLIDE 6: Demo Flow ───────────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Live Demo Flow", "What you're about to see")

steps = [
    ("1", "Submit RITM", "tools/submit_and_ingest.ps1 -Item Azure", "ServiceNow creates a real RITM."),
    ("2", "Intake", "/tickets/intake POST", "gpt-5.4 classifies as azure_infra_change."),
    ("3", "CodeGen", "Terraform PR", "Jinja2 renders disk_iops_change.tf.j2."),
    ("4", "Human gate", "GitHub PR review", "Reviewer merges (the only manual step)."),
    ("5", "Apply", "GH Actions: terraform apply", "Real Azure disk gets new IOPS."),
    ("6", "Verify + Close", "Webhook → Verify → Close", "RITM moves to Closed Complete."),
]
y = Inches(1.3)
for num, title, cmd, body in steps:
    chip = box(s, Inches(0.6), y, Inches(0.6), Inches(0.7), GOLD, shape=MSO_SHAPE.OVAL)
    label(chip, num, color=WHITE, size=20)
    title_box = textbox(s, Inches(1.4), y, Inches(2.5), Inches(0.4),
                       title, font_size=16, color=NAVY, bold=True)
    textbox(s, Inches(1.4), y + Inches(0.36), Inches(4.5), Inches(0.3),
            cmd, font_size=10, color=BLUE)
    textbox(s, Inches(6.0), y + Inches(0.1), Inches(7.0), Inches(0.5),
            body, font_size=14, color=DARK)
    y += Inches(0.85)

textbox(s, Inches(0.6), Inches(6.4), SW - Inches(1.2), Inches(0.5),
        "Watch the dashboard at /api/tickets/<RITM> — stages flip in real time.",
        font_size=14, color=TEAL, bold=True)

notes(s, """
This is exactly what you'll see in the next few minutes.

Step 1: I run a PowerShell script that hits the ServiceNow REST API and creates a real RITM. You'll see the RITM number.

Step 2: The script then POSTs to SnowIaC's /tickets/intake. Intake agent runs, gpt-5.4 classifies it as azure_infra_change.

Step 3: CodeGen renders Terraform from the disk-IOPS template. You'll see the PR open in the GitHub tab.

Step 4: This is the only manual moment. I review the PR — it shows the exact diff — and merge it. That's the human gate, and it's the right place for one.

Step 5: GitHub Actions kicks in. terraform apply against the real managed disk in Azure.

Step 6: The apply finishes, webhook fires back to SnowIaC, Verify agent reads the live disk and confirms IOPS and throughput match the request, Close agent patches the RITM to Closed Complete.

Throughout all of this, the dashboard at /api/tickets/<RITM> shows the stage flipping in real time: INTAKE → CODEGEN → PR OPENED → AWAITING HUMAN MERGE → DEPLOY APPLIED → VERIFIED → COMPLETE.
""")

# ─── SLIDE 7: Why it matters ──────────────────────────────────────────────────
s = add_slide()
header_bar(s, "Why It Matters", "Bigger than disk-resize tickets")

# Big stat cards
cards = [
    ("Hours → seconds", "Ticket-to-PR time", BLUE),
    ("0 → ∞", "Tickets a team can absorb", TEAL),
    ("100%", "Audit trail across SNOW / GH / Azure", NAVY),
    ("Reproducible", "Whole stack in Terraform", GOLD),
]
y = Inches(1.5)
cw = Inches(3.0)
gap = Inches(0.15)
x0 = Inches(0.5)
for i, (big, sub, color) in enumerate(cards):
    x = x0 + i * (cw + gap)
    b = box(s, x, y, cw, Inches(2.0), color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = big
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    sp = tf.add_paragraph()
    sp.alignment = PP_ALIGN.CENTER
    sp.text = sub
    sp.font.size = Pt(13)
    sp.font.color.rgb = WHITE

textbox(s, Inches(0.5), Inches(4.0), SW - Inches(1), Inches(0.5),
        "Same pattern fits any ITSM-to-cloud workflow:", font_size=18, color=NAVY, bold=True)

domains = [
    "VM resize / disk IOPS",
    "Storage account provisioning",
    "Network firewall rule change",
    "RBAC role grants",
    "Key Vault secret rotation",
    "Cost guardrail enforcement",
]
y = Inches(4.7)
for i, item in enumerate(domains):
    row, col = divmod(i, 3)
    x = Inches(0.6) + col * Inches(4.2)
    yy = y + row * Inches(0.5)
    dot = box(s, x, yy, Inches(0.25), Inches(0.25), TEAL, shape=MSO_SHAPE.OVAL)
    textbox(s, x + Inches(0.4), yy - Inches(0.05), Inches(3.7), Inches(0.4),
            item, font_size=14, color=DARK)

notes(s, """
Why this matters beyond one disk-resize demo.

Four numbers to anchor:

— Hours to seconds: the ticket-to-PR time. The agent renders Terraform and opens the PR in well under a minute.

— Zero to infinity: how many tickets a team can absorb. The bottleneck stops being human triage and becomes Azure API throughput.

— 100% audit trail: every state transition is recorded across ServiceNow, GitHub, and Azure activity log. Auditors love this.

— Reproducible: the entire stack is in Terraform, including the Foundry role assignments. We've blown it away and rebuilt it from scratch.

And the bigger point: this isn't really about disk IOPS. The same agent pattern fits every common ITSM-to-cloud workflow — VM resizes, storage account creation, firewall rules, RBAC grants, secret rotation, cost guardrails. Any structured ticket where the change can be expressed as code.

This is the template for autonomous cloud operations.
""")

# ─── SLIDE 8: Closing ─────────────────────────────────────────────────────────
s = add_slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.6), SW, Inches(0.12))
accent.line.fill.background()
accent.fill.solid()
accent.fill.fore_color.rgb = GOLD

textbox(s, Inches(0.8), Inches(2.2), SW - Inches(1.6), Inches(1.0),
        "Let's run it.", font_size=64, color=WHITE, bold=True)
textbox(s, Inches(0.8), Inches(3.3), SW - Inches(1.6), Inches(0.7),
        "Real ServiceNow • Real Azure • Real GitHub PR",
        font_size=22, color=GOLD)
textbox(s, Inches(0.8), Inches(4.0), SW - Inches(1.6), Inches(0.6),
        "github.com/lindazhang2000/iac-ai-agents",
        font_size=16, color=RGBColor(0xCF, 0xDA, 0xE8))

notes(s, """
That's the pitch. Time to switch to the terminal and the browser and run it for real.

I'll have three windows open: the SnowIaC dashboard showing ticket state, the GitHub PR tab, and ServiceNow showing the RITM. Watch all three update as the agents work.

The code is open at github.com/lindazhang2000/iac-ai-agents — Terraform, agent code, GitHub Actions workflow, demo scripts, everything.

Let's go.
""")

prs.save(OUT)
print(f"WROTE {OUT}")
